import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

def apply_to_shape(shape, slide_idx, shape_path, translated_data):
    if shape.has_text_frame:
        for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            for run_idx, run in enumerate(paragraph.runs):
                text = run.text.strip()
                if text:
                    key = f"{slide_idx}_{shape_path}_{p_idx}_{run_idx}"
                    if key in translated_data:
                        run.text = translated_data[key]
                        
    if shape.has_table:
        for row_idx, row in enumerate(shape.table.rows):
            for col_idx, cell in enumerate(row.cells):
                for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        text = run.text.strip()
                        if text:
                            key = f"table_{slide_idx}_{shape_path}_{row_idx}_{col_idx}_{p_idx}_{run_idx}"
                            if key in translated_data:
                                run.text = translated_data[key]
                                
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for idx, child_shape in enumerate(shape.shapes):
            apply_to_shape(child_shape, slide_idx, f"{shape_path}_g{idx}", translated_data)

def apply_text(file_path, json_path, out_path):
    prs = Presentation(file_path)
    with open(json_path, 'r', encoding='utf-8') as f:
        translated_data = json.load(f)
        
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            apply_to_shape(shape, slide_idx, str(shape_idx), translated_data)
            
    prs.save(out_path)
    print(f"Applying done. Saved to {out_path}.")

if __name__ == '__main__':
    apply_text(sys.argv[1], sys.argv[2], sys.argv[3])
