import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_images_from_slides(pptx_path, target_slides, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    prs = Presentation(pptx_path)
    
    for slide_num in target_slides:
        # slide index is 0-based
        if slide_num < 1 or slide_num > len(prs.slides):
            print(f"Slide {slide_num} is out of range.")
            continue
            
        slide = prs.slides[slide_num - 1]
        image_count = 1
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                img_bytes = image.blob
                ext = image.ext
                filename = f"slide_{slide_num}_img_{image_count}.{ext}"
                filepath = os.path.join(output_dir, filename)
                
                with open(filepath, "wb") as f:
                    f.write(img_bytes)
                print(f"Extracted: {filepath}")
                image_count += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # check group shapes recursively if needed
                for child in shape.shapes:
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = child.image
                        img_bytes = image.blob
                        ext = image.ext
                        filename = f"slide_{slide_num}_img_{image_count}_grouped.{ext}"
                        filepath = os.path.join(output_dir, filename)
                        
                        with open(filepath, "wb") as f:
                            f.write(img_bytes)
                        print(f"Extracted: {filepath}")
                        image_count += 1

if __name__ == "__main__":
    pptx_file = r"C:\Users\tetsu\Desktop\尾瀬デュエル：地域課題を「連鎖」で解き明かす振り返りセッション デザイン変更_zh-TW.pptx"
    out_dir = r"d:\★Antigravity\AudioBlogProject\extracted_images"
    slides_to_extract = [1, 3, 4, 5, 6, 23]
    
    extract_images_from_slides(pptx_file, slides_to_extract, out_dir)
    print("Done extracting images.")
