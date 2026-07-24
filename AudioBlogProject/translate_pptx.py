import sys
import os
import collections
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from deep_translator import GoogleTranslator

# Using deep-translator for reliable translation
translator = GoogleTranslator(source='ja', target='zh-TW')

def translate_text(text):
    if not text or not text.strip():
        return text
    try:
        # Translate blocks
        return translator.translate(text)
    except Exception as e:
        print(f"Translation Error for context {text}: {e}")
        return text

def translate_paragraph(paragraph):
    # Paragraphs contain 1 or more runs. Sentences can be split across runs.
    # To get a coherent translation, we combine the text, translate it,
    # and put it all in the first run, preserving the first run's font styling.
    # We then clear the subsequent runs.
    
    combined_text = "".join(run.text for run in paragraph.runs)
    if not combined_text.strip():
        return
    
    translated_text = translate_text(combined_text)
    
    # Write back
    if paragraph.runs:
        paragraph.runs[0].text = translated_text
        # Optional: ensure we preserve subsequent runs as empty
        for i in range(1, len(paragraph.runs)):
            paragraph.runs[i].text = ""

def process_shape(shape):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            translate_paragraph(paragraph)
            
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    translate_paragraph(paragraph)
                    
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            process_shape(child)

def main():
    input_file = r"C:\Users\tetsu\Desktop\尾瀬デュエル：地域課題を「連鎖」で解き明かす振り返りセッション デザイン変更.pptx"
    output_file = r"C:\Users\tetsu\Desktop\尾瀬デュエル：地域課題を「連鎖」で解き明かす振り返りセッション デザイン変更_zh-TW.pptx"

    if not os.path.exists(input_file):
        print(f"File not found: {input_file}")
        sys.exit(1)

    print(f"Loading presentation: {input_file}")
    prs = Presentation(input_file)

    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        print(f"Processing slide {i+1} / {total_slides}...")
        for shape in slide.shapes:
            process_shape(shape)
            
    try:
        print(f"Saving to {output_file}")
        prs.save(output_file)
        print("Done!")
    except Exception as e:
        print(f"Failed to save: {e}")

if __name__ == '__main__':
    main()
