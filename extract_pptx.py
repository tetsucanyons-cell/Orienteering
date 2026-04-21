import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

def extract_from_shape(shape, slide_idx, shape_path, extracted_data):
    if shape.has_text_frame:
        for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
            for run_idx, run in enumerate(paragraph.runs):
                text = run.text.strip()
                if text:
                    key = f"{slide_idx}_{shape_path}_{p_idx}_{run_idx}"
                    extracted_data[key] = text
                    
    if shape.has_table:
        for row_idx, row in enumerate(shape.table.rows):
            for col_idx, cell in enumerate(row.cells):
                for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        text = run.text.strip()
                        if text:
                            key = f"table_{slide_idx}_{shape_path}_{row_idx}_{col_idx}_{p_idx}_{run_idx}"
                            extracted_data[key] = text
                            
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for idx, child_shape in enumerate(shape.shapes):
            extract_from_shape(child_shape, slide_idx, f"{shape_path}_g{idx}", extracted_data)

def extract_text(file_path):
    prs = Presentation(file_path)
    extracted_data = {}
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            extract_from_shape(shape, slide_idx, str(shape_idx), extracted_data)
            
    with open('d:/★Antigravity/extracted_texts.json', 'w', encoding='utf-8') as f:
        json.dump(extracted_data, f, ensure_ascii=False, indent=2)
        
    print(f"Extraction done. {len(extracted_data)} runs found.")

if __name__ == '__main__':
    extract_text(sys.argv[1])
